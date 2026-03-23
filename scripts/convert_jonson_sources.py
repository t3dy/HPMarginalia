"""Convert all documents in 'Ben Jonson Life/' to markdown files.

Handles: PDF (PyMuPDF), PPTX (python-pptx), EPUB (zipfile+bs4), TXT (copy).
Outputs to 'Ben Jonson Life/md/' with clean filenames.
"""

import sys
import re
import shutil
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent
SOURCE_DIR = BASE_DIR / "Ben Jonson Life"
OUTPUT_DIR = SOURCE_DIR / "md"


def slugify_filename(name: str) -> str:
    """Create a short, readable filename slug."""
    # Remove file extension
    name = Path(name).stem
    # Remove common noise patterns
    name = re.sub(r'\{[^}]*\}', '', name)  # {metadata}
    name = re.sub(r'\[[^\]]*\]', '', name)  # [series info]
    name = re.sub(r'libgen[.\s]li', '', name, flags=re.IGNORECASE)
    name = re.sub(r'\d{5,}', '', name)  # long numbers
    name = name.strip(' -_.,')
    # Truncate to something reasonable
    words = name.split()
    if len(words) > 8:
        words = words[:8]
    slug = '_'.join(words)
    slug = re.sub(r'[^\w\s-]', '', slug)
    slug = re.sub(r'[\s]+', '_', slug)
    slug = re.sub(r'_+', '_', slug)
    slug = slug.strip('_')
    return slug[:80]


def convert_pdf(pdf_path: Path, md_path: Path):
    """Extract text from PDF using PyMuPDF, write as markdown."""
    import fitz
    doc = fitz.open(str(pdf_path))
    page_count = len(doc)
    lines = []
    lines.append(f"# {pdf_path.stem[:100]}\n")
    lines.append(f"*Source: {pdf_path.name}*\n")
    lines.append(f"*Pages: {page_count}*\n\n---\n")

    for page_num in range(page_count):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        if text.strip():
            lines.append(f"\n## Page {page_num + 1}\n")
            lines.append(text)
    doc.close()

    md_path.write_text('\n'.join(lines), encoding='utf-8')
    print(f"  PDF -> MD: {pdf_path.name} ({page_count} pages)")


def convert_pptx(pptx_path: Path, md_path: Path):
    """Extract text from PPTX using python-pptx, write as markdown."""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    lines = []
    lines.append(f"# {pptx_path.stem[:100]}\n")
    lines.append(f"*Source: {pptx_path.name}*\n")
    lines.append(f"*Slides: {len(prs.slides)}*\n\n---\n")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Slide {i}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
                lines.append("")
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                lines.append("")

    md_path.write_text('\n'.join(lines), encoding='utf-8')
    print(f"  PPTX -> MD: {pptx_path.name} ({len(prs.slides)} slides)")


def convert_epub(epub_path: Path, md_path: Path):
    """Extract text from EPUB using zipfile + BeautifulSoup."""
    import zipfile
    from bs4 import BeautifulSoup

    lines = []
    lines.append(f"# {epub_path.stem[:100]}\n")
    lines.append(f"*Source: {epub_path.name}*\n\n---\n")

    chapter_count = 0
    with zipfile.ZipFile(str(epub_path), 'r') as zf:
        # Find content files (XHTML/HTML)
        content_files = sorted([
            n for n in zf.namelist()
            if n.endswith(('.xhtml', '.html', '.htm'))
            and 'toc' not in n.lower()
            and 'nav' not in n.lower()
        ])

        for cf in content_files:
            try:
                html = zf.read(cf).decode('utf-8', errors='replace')
                soup = BeautifulSoup(html, 'html.parser')
                # Remove scripts and styles
                for tag in soup(['script', 'style']):
                    tag.decompose()
                text = soup.get_text(separator='\n')
                text = re.sub(r'\n{3,}', '\n\n', text).strip()
                if text and len(text) > 20:
                    chapter_count += 1
                    lines.append(f"\n## Chapter/Section {chapter_count}\n")
                    lines.append(text)
            except Exception as e:
                lines.append(f"\n*[Error reading {cf}: {e}]*\n")

    md_path.write_text('\n'.join(lines), encoding='utf-8')
    print(f"  EPUB -> MD: {epub_path.name} ({chapter_count} sections)")


def convert_txt(txt_path: Path, md_path: Path):
    """Wrap existing text file in minimal markdown."""
    text = txt_path.read_text(encoding='utf-8', errors='replace')
    # Convert --- PAGE N --- markers to ## Page N headers
    text = re.sub(r'^---\s*PAGE\s+(\d+)\s*---\s*$', r'## Page \1', text, flags=re.MULTILINE)
    lines = []
    lines.append(f"# {txt_path.stem[:100]}\n")
    lines.append(f"*Source: {txt_path.name}*\n\n---\n")
    lines.append(text)

    md_path.write_text('\n'.join(lines), encoding='utf-8')
    print(f"  TXT -> MD: {txt_path.name}")


def main():
    OUTPUT_DIR.mkdir(exist_ok=True)

    files = list(SOURCE_DIR.iterdir())
    convertible = [f for f in files if f.is_file() and f.suffix.lower() in ('.pdf', '.pptx', '.epub', '.txt')]

    if not convertible:
        print("No convertible files found in", SOURCE_DIR)
        return

    print(f"Converting {len(convertible)} files from {SOURCE_DIR.name}/\n")

    converters = {
        '.pdf': convert_pdf,
        '.pptx': convert_pptx,
        '.epub': convert_epub,
        '.txt': convert_txt,
    }

    results = []
    for f in sorted(convertible):
        slug = slugify_filename(f.name)
        md_path = OUTPUT_DIR / f"{slug}.md"
        converter = converters.get(f.suffix.lower())
        if converter:
            try:
                converter(f, md_path)
                results.append((f.name, md_path.name, "OK"))
            except Exception as e:
                print(f"  ERROR: {f.name}: {e}")
                results.append((f.name, md_path.name, f"ERROR: {e}"))

    print(f"\n{'='*60}")
    print(f"Converted {sum(1 for _, _, s in results if s == 'OK')}/{len(results)} files")
    print(f"Output directory: {OUTPUT_DIR}")
    print()
    for src, dst, status in results:
        print(f"  {status:5s}  {dst}")


if __name__ == "__main__":
    main()
